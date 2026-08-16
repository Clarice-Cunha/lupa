"""Preenche o template oficial da apresentação da Etapa Regional.

O template do HackaNAV tem 9 slides, dos quais 5 precisam ser preenchidos.
Este script escreve o conteúdo neles sem quebrar a formatação: em vez de criar
parágrafos novos do zero, ele copia um parágrafo existente e troca só o texto.
Assim fonte, tamanho e cor do tema continuam valendo.

Uso (a partir da raiz do projeto):
    python docs/gerar_apresentacao_regional.py
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches

RAIZ = Path(__file__).resolve().parent.parent
MODELO = Path.home() / "Downloads" / "Hacka2026 - Template Apresentação Etapa Regional.pptx"
DESTINO = RAIZ / "docs" / "APRESENTACAO_REGIONAL_LUPA.pptx"
PASTA_IMAGENS = RAIZ / "docs" / "imagens"

# ---------------------------------------------------------------- dados fixos

EQUIPE = "CB"
ESCOLA = "Contemporâneo Lagoa Nova"
PROJETO = "LUPA — Leitor de URLs, Plataformas e Audiovisuais"
PROFESSOR = "Hector Gabriel Ribeiro Liberalino"

# ATENÇÃO: o Portal do Nave a Vela registra "Miguel Cavalcanti Filgueira" como
# terceiro integrante, mas quem participa e vai falar na live é o Pedro. Isso
# precisa ser resolvido com o professor antes do envio — a apresentação mostra
# quem realmente está na equipe.
ALUNOS = [
    "Clarice Cunha Pinto (capitã)",
    "Benjamim de Almeida das Chagas",
    "Pedro Moreno de Lima Bessa",
]

# Cada item é (texto, negrito). Texto vazio vira linha em branco.
#
# O texto do slide 7 é curto de propósito: quem conta o erro é o fluxograma
# logo abaixo, e ele precisa de espaço para continuar legível na transmissão.
# A explicação completa fica na fala, não na tela.
ERRO = [
    ("O desafio: dar mais certeza à nota", True),
    ("A pontuação vinha de sinais indiretos — conexão segura, idade do "
     "domínio, linguagem sensacionalista. Nenhum deles diz se o conteúdo "
     "é verdadeiro.", False),
    ("", False),
    ("A decisão: em 1º de maio, às 14h23, integramos o banco de checagens "
     "do IFCN.", True),
]

# A coluna da esquerda do slide 8 é estreita e o logo do HackaNAV ocupa o
# canto inferior. O texto precisa caber em cerca de 13 linhas.
APRENDIZADO = [
    ("A reação", True),
    ("Frustração — cogitamos abandonar a API. Discordamos da sugestão de "
     "remover a análise e procuramos outra saída.", False),
    ("", False),
    ("A correção, no mesmo dia", True),
    ("17h11 — não consultar o banco em páginas iniciais", False),
    ("17h21 — filtro descarta checagens sem relação real", False),
    ("", False),
    ("O que aprendemos", True),
    ("Tratamos um sinal de alerta como se fosse uma prova — o mesmo erro que "
     "a desinformação explora nas pessoas.", False),
    ("", False),
    ("Oito dias depois: nota 100 passou a vir com um aviso.", False),
]

# ------------------------------------------------------------------- funções


def achar_shape(slide, nome: str):
    for shape in slide.shapes:
        if shape.name == nome:
            return shape
    raise KeyError(f"shape {nome!r} nao encontrado")


def texto_do_run(paragrafo_xml, texto: str) -> None:
    """Deixa o parágrafo com um único run contendo o texto informado."""
    runs = paragrafo_xml.findall(qn("a:r"))
    for extra in runs[1:]:
        paragrafo_xml.remove(extra)
    if not runs:
        return
    runs[0].find(qn("a:t")).text = texto


def definir_negrito(paragrafo_xml, negrito: bool) -> None:
    run = paragrafo_xml.find(qn("a:r"))
    if run is None:
        return
    rpr = run.find(qn("a:rPr"))
    if rpr is None:
        rpr = run.makeelement(qn("a:rPr"), {})
        run.insert(0, rpr)
    rpr.set("b", "1" if negrito else "0")


def escrever(text_frame, itens, modelo_idx: int = 0, alinhamento=PP_ALIGN.LEFT):
    """Reescreve o quadro de texto usando um parágrafo existente como molde.

    Copiar o XML do parágrafo preserva a fonte e a cor do tema, que não são
    recuperáveis pela API quando se cria um parágrafo do zero.
    """
    corpo = text_frame._txBody
    molde = copy.deepcopy(text_frame.paragraphs[modelo_idx]._p)

    for paragrafo in corpo.findall(qn("a:p")):
        corpo.remove(paragrafo)

    for texto, negrito in itens:
        novo = copy.deepcopy(molde)
        texto_do_run(novo, texto)
        definir_negrito(novo, negrito)
        corpo.append(novo)

    for paragrafo in text_frame.paragraphs:
        paragrafo.alignment = alinhamento


def main() -> None:
    if not MODELO.exists():
        raise SystemExit(f"Template nao encontrado em {MODELO}")

    pr = Presentation(str(MODELO))
    slides = pr.slides

    # Slide 3 — identificação da equipe
    escrever(achar_shape(slides[2], "Google Shape;1197;p7").text_frame,
             [(EQUIPE, True)], alinhamento=PP_ALIGN.CENTER)
    escrever(achar_shape(slides[2], "Google Shape;1198;p7").text_frame,
             [(ESCOLA, False)], alinhamento=PP_ALIGN.CENTER)

    # Slide 4 — integrantes. O molde é o parágrafo 1, que não é negrito.
    integrantes = [("Alunos:", True)]
    integrantes += [(nome, False) for nome in ALUNOS]
    integrantes += [("", False), ("Professor(a):", True), (PROFESSOR, False)]
    escrever(achar_shape(slides[3], "Google Shape;1323;p8").text_frame,
             integrantes, modelo_idx=1)

    # Slide 5 — nome do projeto
    escrever(achar_shape(slides[4], "Google Shape;1860;g3d1cff3b8f6_2_1917").text_frame,
             [(PROJETO, True)], alinhamento=PP_ALIGN.CENTER)

    # Slides 7 e 8 — a narrativa do erro, em duas partes
    escrever(achar_shape(slides[6], "Google Shape;2610;p15").text_frame, ERRO)

    caixa_aprendizado = achar_shape(slides[7], "Google Shape;2997;g3c919147bba_1_18")
    # Sobe a caixa para logo abaixo do título: o texto cresce para baixo e,
    # na posição original, a última linha ficava atrás do logo do HackaNAV.
    caixa_aprendizado.top = Inches(1.50)
    escrever(caixa_aprendizado.text_frame, APRENDIZADO)

    inserir_imagens(slides)

    pr.save(str(DESTINO))
    print(f"Gerado: {DESTINO.relative_to(RAIZ)}")
    print()
    print("Falta fazer a mao:")
    print("  - slide 6: inserir o video da etapa anterior (Inserir > Video)")
    print("  - conferir o nome do 3o integrante no Portal (Miguel x Pedro)")


def inserir_imagens(slides) -> None:
    """Coloca os fluxogramas nos dois slides da narrativa do erro.

    As medidas foram tiradas das formas decorativas do template: no slide 7
    sobra uma faixa larga abaixo do texto, e no slide 8 há um quadro vazio
    do lado direito. Em ambos os casos a imagem é encaixada sem cobrir o
    logo do HackaNAV nem a moldura.
    """
    faixa_slide7 = PASTA_IMAGENS / "falso-positivo.png"
    quadro_slide8 = PASTA_IMAGENS / "cinco-camadas.png"

    for caminho in (faixa_slide7, quadro_slide8):
        if not caminho.exists():
            raise SystemExit(
                f"Falta {caminho.name}. Rode antes: "
                "python docs/gerar_imagens_fluxogramas.py"
            )

    # Slide 7 — faixa livre entre o fim do texto e o rodapé da moldura
    largura = 5.70
    slides[6].shapes.add_picture(
        str(faixa_slide7),
        Inches(4.965 - largura / 2), Inches(2.55),
        width=Inches(largura),
    )

    # Slide 8 — quadro vazio da direita (x 5.04 a 9.53, y 0.70 a 5.19)
    largura = 4.00
    altura = largura * 1120 / 1560
    slides[7].shapes.add_picture(
        str(quadro_slide8),
        Inches(5.04 + (4.49 - largura) / 2),
        Inches(0.70 + (4.49 - altura) / 2),
        width=Inches(largura),
    )


if __name__ == "__main__":
    main()
